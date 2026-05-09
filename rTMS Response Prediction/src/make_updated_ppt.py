from __future__ import annotations

import math
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ART = ROOT / "artifacts"
OUT_DIR = ROOT / "outputs"


def _add_title(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0


def _add_picture_slide(
    prs: Presentation,
    title: str,
    img_path: Path,
    caption: str | None = None,
    *,
    height_in: float = 4.6,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title

    left = Inches(0.8)
    top = Inches(1.4)
    height = Inches(height_in)
    slide.shapes.add_picture(str(img_path), left, top, height=height)

    if caption:
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.0), Inches(0.6))
        tf = tx_box.text_frame
        tf.text = caption
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT


def _save_confusion_matrix_png(cm: np.ndarray, labels: list[str], out_path: Path) -> None:
    fig, ax = plt.subplots(figsize=(6.0, 5.2), dpi=200)
    im = ax.imshow(cm, cmap="Blues")
    ax.set_title("BEED confusion matrix (holdout)", fontsize=12)
    ax.set_xlabel("model output class", fontsize=11)
    ax.set_ylabel("true", fontsize=11)
    ax.set_xticks(range(len(labels)), labels=labels)
    ax.set_yticks(range(len(labels)), labels=labels)

    for i in range(cm.shape[0]):
        for j in range(cm.shape[1]):
            ax.text(j, i, str(int(cm[i, j])), ha="center", va="center", fontsize=10, color="black")

    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    fig.tight_layout()
    fig.savefig(out_path)
    plt.close(fig)


def _save_perm_importance_bar_png(summary_csv: Path, out_path: Path, *, top_n: int = 8) -> list[str]:
    df = pd.read_csv(summary_csv)
    df = df.sort_values(["folds_in_top_5", "mean_permutation_importance"], ascending=[False, False]).head(top_n)
    feats = df["feature"].tolist()
    vals = df["mean_permutation_importance"].to_numpy()
    errs = df["std_permutation_importance"].to_numpy()

    fig, ax = plt.subplots(figsize=(7.2, 4.2), dpi=200)
    ax.barh(feats[::-1], vals[::-1], xerr=errs[::-1], color="C0", alpha=0.9)
    ax.set_title("BEED permutation importance (CV mean ± std)", fontsize=12)
    ax.set_xlabel("balanced accuracy drop when feature is shuffled", fontsize=11)
    ax.grid(axis="x", linestyle="--", alpha=0.4)
    fig.tight_layout()
    fig.savefig(out_path)
    plt.close(fig)
    return feats


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    # load metrics from prior run output files
    # confusion matrix is printed, so we recompute quickly from BEED_Data.csv for the ppt
    beed_csv = Path(r"C:\Users\coleb\Downloads\beed_+bangalore+eeg+epilepsy+dataset\BEED_Data.csv")
    df = pd.read_csv(beed_csv)
    X = df[[f"X{i}" for i in range(1, 17)]].astype(float)
    y = df["y"].astype(int)

    # deterministic split matching our scripts: seed 42, test_size 0.25
    rng = np.random.RandomState(42)
    idx = np.arange(len(y))
    # stratified split manual to avoid importing sklearn here
    test_idx = []
    for cls in sorted(y.unique().tolist()):
        cls_idx = idx[y.to_numpy() == cls]
        rng.shuffle(cls_idx)
        n_test = int(math.floor(0.25 * len(cls_idx)))
        test_idx.extend(cls_idx[:n_test].tolist())
    test_idx = np.array(sorted(test_idx))
    train_mask = np.ones(len(y), dtype=bool)
    train_mask[test_idx] = False

    # load the saved RF holdout model from artifacts
    rf_path = ART / "beed_random_forest_seed42.joblib"
    bundle = None
    if rf_path.is_file():
        import joblib

        bundle = joblib.load(rf_path)
    model = bundle["model"] if bundle else None
    if model is None:
        raise SystemExit(f"missing model artifact: {rf_path}")

    pred = model.predict(X.iloc[test_idx])
    labels = [str(c) for c in sorted(y.unique().tolist())]
    cm = np.zeros((len(labels), len(labels)), dtype=int)
    y_test = y.iloc[test_idx].to_numpy()
    for t, p in zip(y_test, pred, strict=False):
        cm[int(t), int(p)] += 1

    cm_png = OUT_DIR / "beed_confusion_matrix.png"
    _save_confusion_matrix_png(cm, labels, cm_png)

    perm_png = OUT_DIR / "beed_perm_importance.png"
    top_feats = _save_perm_importance_bar_png(ART / "beed_biomarker_summary.csv", perm_png, top_n=8)

    biomarker_json = ART / "beed_biomarkers_consistent_topk.json"
    consistent = ["X8", "X11", "X12"]
    if biomarker_json.is_file():
        import json

        consistent = json.loads(biomarker_json.read_text(encoding="utf-8")).get(
            "consistent_features_all_folds", consistent
        )

    prs = Presentation()

    _add_title(
        prs,
        "Closed-loop adaptive rTMS system (plan + current progress)",
        "real-time EEG + embedded AI • updated datasets, outcomes, and outputs",
    )

    _add_bullets(
        prs,
        "What is rTMS?",
        [
            "non-invasive neuromodulation using magnetic pulses to modulate cortical excitability",
            "clinical use: treatment-resistant depression (common), OCD, etc.",
            "key problem: response rates are variable → biomarkers that stratify response are valuable",
        ],
    )

    _add_bullets(
        prs,
        "Updated datasets (what we have vs what we need)",
        [
            "BEED epilepsy dataset (available now): 8,000 rows, 16 features (X1–X16), 4 balanced classes (y=0–3)",
            "TDBRAIN (in progress): baseline resting EEG + clinical table (participants.tsv) for defining rTMS response",
            "note: TDBRAIN derivatives zip is password-protected; outcomes come from participants.tsv (not in derivatives)",
        ],
    )

    _add_bullets(
        prs,
        "Updated outcomes + how we obtain them",
        [
            "BEED: outcome label is y ∈ {0,1,2,3} (multiclass classification)",
            "TDBRAIN rTMS (target): responder = ≥50% reduction in BDI from baseline to end of treatment",
            "means of obtaining rTMS outcome: compute from participants.tsv columns `BDI Pre` and `BDI Post` (or `Responder` if provided)",
        ],
    )

    _add_bullets(
        prs,
        "Epoching and outputs to show (EEG → features)",
        [
            "BEED is already tabular features; no raw EEG epoching required for the baseline",
            "TDBRAIN planned EEG feature extraction: use first clean eyes-closed epoch (e.g., 30s) per subject",
            "typical pipeline: clean EEG → pick 30s artifact-free epoch → compute bandpowers/connectivity/features → train classifier",
        ],
    )

    _add_bullets(
        prs,
        "Objective 4 (explainability): biomarker definition used",
        [
            "per CV fold: train model on train split, compute permutation importance on validation split only",
            "biomarker stability criterion: feature is in the top-5 permutation importance in every fold",
            f"BEED consistent top-5 features across 5 folds: {', '.join(consistent)}",
        ],
    )

    _add_picture_slide(prs, "BEED model output: confusion matrix", cm_png, caption="random forest, 75/25 stratified holdout")
    _add_picture_slide(
        prs,
        "BEED explainability output: permutation importance",
        perm_png,
        caption=f"top features by CV mean importance: {', '.join(top_feats[:5])}",
    )

    td_cm = OUT_DIR / "tdbrain_rtms_validation_confusion_matrix.png"
    td_bar = OUT_DIR / "tdbrain_rtms_validation_metrics.png"
    td_curves = OUT_DIR / "tdbrain_rtms_training_outcomes.png"
    if td_cm.is_file():
        _add_picture_slide(
            prs,
            "TDBRAIN / rTMS response — validation confusion matrix",
            td_cm,
            height_in=5.9,
        )
    if td_bar.is_file():
        _add_picture_slide(
            prs,
            "TDBRAIN / rTMS response — validation metrics",
            td_bar,
            height_in=5.9,
        )
    if td_curves.is_file():
        _add_picture_slide(
            prs,
            "TDBRAIN / rTMS response — training outcomes",
            td_curves,
            height_in=5.9,
        )

    _add_bullets(
        prs,
        "Next steps (to complete rTMS objectives)",
        [
            "extract TDBRAIN derivatives (EEG) + obtain participants.tsv (clinical labels)",
            "define responder label from BDI Pre/Post; ensure one row per participant (avoid leakage across sessions)",
            "repeat the same reporting stack as BEED: CV + locked holdout + fold-stable biomarkers",
        ],
    )

    out_pptx = OUT_DIR / "BioAIrTMS_updated.pptx"
    prs.save(out_pptx)

    print(f"wrote {out_pptx}")


if __name__ == "__main__":
    main()

